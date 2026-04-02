import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = '/Users/hengyuhu/mainvoice-ai'

# --- Constants & Colors ---
BG_COLOR = RGBColor(250, 248, 244) # #FAF8F4 (Warm Neutral)
TEXT_DARK = RGBColor(45, 50, 32)   # #2D3220 (Dark Sage)
TEXT_MUTED = RGBColor(106, 117, 96)# #6A7560 (Muted Sage)
ACCENT_GREEN = RGBColor(107, 143, 113) # #6B8F71 (Sage Green)
ACCENT_CORAL = RGBColor(232, 133, 106) # #E8856A (Soft Coral)

FONT_SERIF = "Georgia"      # Fallback for Fraunces
FONT_SANS = "Arial"         # Fallback for DM Sans

prs = Presentation()
# Set presentation size to widescreen (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Clear default slides, use blank layout (index 6 is usually blank)
blank_slide_layout = prs.slide_layouts[6] 

def apply_background(slide):
    # Add a full-screen rectangle for background color
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_COLOR
    shape.line.fill.background()
    # Send to back by manipulating underlying xml
    slide.shapes._spTree.insert(2, shape._element)

def add_title(slide, text, top=Inches(0.8), left=Inches(1.0), width=Inches(11.3), font_size=44, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.name = FONT_SERIF
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_body(slide, text, top=Inches(2.0), left=Inches(1.0), width=Inches(5.0), font_size=24, color=TEXT_MUTED, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_SANS
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_image(slide, image_path, left=Inches(6.5), top=Inches(1.5), width=Inches(6.0), height=None):
    if os.path.exists(image_path):
        if height:
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        else:
            slide.shapes.add_picture(image_path, left, top, width=width)
    else:
        print(f"Warning: {image_path} not found.")

def create_slide_with_image(title, body_text, image_name):
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide)
    add_title(slide, title)
    add_body(slide, body_text)
    
    img_path = os.path.join("mockups", image_name)
    add_image(slide, img_path)
    
    # Add gentle accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(1.5), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CORAL
    line.line.fill.background()

# --- Build Slides ---

# Slide 1: Title
s1 = prs.slides.add_slide(blank_slide_layout)
apply_background(s1)
add_title(s1, "MainVoice AI", top=Inches(2.8), font_size=60, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_body(s1, "Your Local, Offline Mental Health Companion\nEmpathetic • Private • Efficacious", top=Inches(4.0), left=Inches(1.0), width=Inches(11.3), font_size=28, align=PP_ALIGN.CENTER)

# Slide 2: The Challenge
s2 = prs.slides.add_slide(blank_slide_layout)
apply_background(s2)
add_title(s2, "The Mental Health Crisis")
body_s2 = ("• Anxiety & Depression among youth have risen dramatically.\n"
          "• High barriers to accessing immediate professional help.\n"
          "• Privacy concerns prevent many from sharing deeply.\n\n"
          "We need a safe, accessible, and scientifically grounded first-response tool.")
add_body(s2, body_s2, width=Inches(10.0))

# Slide 3: Our Solution
s3 = prs.slides.add_slide(blank_slide_layout)
apply_background(s3)
add_title(s3, "Our Solution")
body_s3 = ("MainVoice AI bridges the gap between isolation and therapy.\n\n"
          "• 100% Local & Offline Privacy\n"
          "• CBT-Informed Empathy Engine\n"
          "• Multi-dimensional support ecosystem")
add_body(s3, body_s3, width=Inches(10.0))

# Slide 4: Hero Landing
create_slide_with_image(
    "A Calming Landing Experience", 
    "Designed entirely with therapeutic aesthetics.\n\nWarm Sage Green palettes, rounded glassmorphism, and completely distraction-free environments help users feel safe.", 
    "shot_hero.png"
)

# Slide 5: Core Feature - Chat
create_slide_with_image(
    "Empathetic AI Chat", 
    "Powered by an offline logic engine.\n\nThe AI dynamically responds to sensitive topics without ever sending data to the cloud. \nIt employs CBT re-framing techniques gently.", 
    "shot_chat.png"
)

# Slide 6: Mood Tracker
create_slide_with_image(
    "Mood Tracker & Insights", 
    "Daily check-ins help users articulate ambiguous feelings.\n\nThe system charts these over time to provide Weekly Emotional Reports, offering tangible insights.", 
    "shot_mood.png"
)

# Slide 7: Gratitude Board
create_slide_with_image(
    "Gratitude Board", 
    "A digital sanctuary for highlighting the good.\n\nUsing positive psychology, the AI provides a warm, third-person reflection on the user's daily gratitude entries.", 
    "shot_gratitude.png"
)

# Slide 8: Micro Habits
create_slide_with_image(
    "Daily Micro-Habits", 
    "Rebuilding momentum through tiny wins.\n\nUsers select small, achievable goals (like 'Drink Water' or 'Walk 5 mins') rather than overwhelming macro-goals.", 
    "shot_habits.png"
)

# Slide 9: Parent Letters
create_slide_with_image(
    "Parent Communication Templates", 
    "Bridging generational gaps.\n\nA safe interface for teens to generate structured, emotionally mature letters to their parents or teachers when speaking is too difficult.", 
    "shot_letter.png"
)

# Slide 10: CBT Exercises
create_slide_with_image(
    "Self-Help CBT Manual", 
    "Actionable, evidence-based coping strategies.\n\nStep-by-step breathing guides, grounding techniques, and emotional articulation worksheets immediately accessible.", 
    "shot_manual.png"
)

# Slide 11: Stress Testing
create_slide_with_image(
    "Stress Assessments", 
    "Life & Academic stress quizzes.\n\nProviding immediate, non-diagnostic feedback and helping users validate their feelings before deciding if they need robust professional help.", 
    "shot_quiz.png"
)

# Slide 12: Technology & Privacy
s12 = prs.slides.add_slide(blank_slide_layout)
apply_background(s12)
add_title(s12, "Technology & Privacy-First")
body_s12 = ("• Zero-Tracking Policy\n"
           "• All mental health logs and chat histories are securely stored via local storage within the browser.\n"
           "• No databases, no telemetry, complete user autonomy.\n"
           "• Powered by seamlessly integrated AI APIs with restricted referrer access.")
add_body(s12, body_s12, width=Inches(10.0))

# Slide 13: Crisis Support
create_slide_with_image(
    "Immediate Crisis Support", 
    "The system continuously scans for high-risk safety keywords.\n\nInstantly halts normal conversation to provide direct, non-judgmental emergency numbers and local crisis centers.", 
    "shot_crisis.png"
)

# Slide 14: Future Vision
s14 = prs.slides.add_slide(blank_slide_layout)
apply_background(s14)
add_title(s14, "Future Vision & Impact")
body_s14 = ("MainVoice AI represents the future of democratized, safe mental health triage.\n\n"
           "Our vision is deeply rooted in providing immediate emotional first-aid that compliments, rather than replaces, traditional therapy.")
add_body(s14, body_s14, width=Inches(10.0))

# Slide 15: Thank You
s15 = prs.slides.add_slide(blank_slide_layout)
apply_background(s15)
add_title(s15, "Thank You", top=Inches(3.0), font_size=60, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_body(s15, "Questions & Answers", top=Inches(4.2), left=Inches(1.0), width=Inches(11.3), font_size=28, align=PP_ALIGN.CENTER)

# Save the presentation
output_path = os.path.join(BASE_DIR, "MainVoice-AI-Presentation.pptx")
prs.save(output_path)
print(f"Presentation generated at {output_path}")
